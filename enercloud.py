from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
DARK_BLUE = RGBColor(0x0D, 0x1B, 0x2A)
MED_BLUE = RGBColor(0x2D, 0x6C, 0xB4)
LIGHT_BLUE = RGBColor(0x4A, 0x9B, 0xD9)
YELLOW = RGBColor(0xF5, 0xB0, 0x41)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x1C, 0x28, 0x33)
GRAY = RGBColor(0x88, 0x88, 0x88)
LIGHT_BG = RGBColor(0xF0, 0xF4, 0xF8)

slides_data = [
    {
        "title": "Cloud Computing for Energy Systems",
        "subtitle": "Software-Defined Grid  |  EnerCloud Pvt Limited",
        "bullets": [],
        "funny": "Electricity is getting smarter",
        "visual": "Ultra clean 3D isometric infographic, blue and yellow theme, soft lighting, corporate professional, no text — futuristic smart city with glowing cloud brain floating above connected buildings, power lines, and solar panels, energy particles flowing upward into cloud",
        "morph": "MORPH: Zoom into city — camera pushes forward into the streets between buildings",
        "notes": "Welcome everyone. Today we explore how cloud computing is transforming energy systems. Think of it as giving the electricity grid a brain.",
        "is_title": True
    },
    {
        "title": "The Real Problem",
        "subtitle": "",
        "bullets": ["Solar gives maximum energy at noon", "People use maximum energy at 6-9 PM", "This mismatch = wasted energy + blackouts", "We need something to bridge this gap"],
        "funny": "Sun works 9 to 5... humans don't",
        "visual": "Ultra clean 3D isometric infographic, blue and yellow theme, soft lighting, corporate professional, no text — split scene showing bright sunny city on left with solar panels glowing and dark evening city on right with homes lit up, visible energy gap between them",
        "morph": "MORPH: Split screen slides apart — left and right sides separate creating space",
        "notes": "This is the fundamental problem. The sun peaks at noon. But humans come home at 6 PM and turn on AC, lights, TV. That gap is where energy gets wasted or blackouts happen.",
        "is_title": False
    },
    {
        "title": "Types of Energy Systems",
        "subtitle": "",
        "bullets": ["Thermal — coal, gas, nuclear power plants", "Renewable — solar farms, wind turbines", "Distributed — rooftop solar, EV chargers, home batteries", "Each type produces and behaves differently"],
        "funny": "Each one has its own mood",
        "visual": "Ultra clean 3D isometric infographic, blue and yellow theme, soft lighting, corporate professional, no text — three distinct 3D icons grouped together: a coal power plant with smokestacks, a solar panel array with sun rays, and a wind turbine with spinning blades",
        "morph": "MORPH: Icons expand — each of the three icons scales up and spreads outward",
        "notes": "We have three main types. Thermal is the old reliable but polluting. Renewable is clean but unpredictable. Distributed is the newest — energy generated right where it's used.",
        "is_title": False
    },
    {
        "title": "The Traditional Grid",
        "subtitle": "Built 100+ years ago",
        "bullets": ["Electricity flows ONE way — plant to home", "One big central power plant serves everyone", "No real-time monitoring or control", "Simple but outdated for today's needs"],
        "funny": "Old system = simple but stuck in the past",
        "visual": "Ultra clean 3D isometric infographic, blue and yellow theme, soft lighting, corporate professional, no text — single large power plant on left with one-way arrows flowing through transmission towers to a row of simple homes on right, clean linear layout",
        "morph": "MORPH: Fade to modern — old power plant dissolves and reforms into multiple smaller sources",
        "notes": "The traditional grid was designed for one purpose — move electricity from a big coal plant to your home. One direction. No thinking. It worked fine when we only had power plants.",
        "is_title": False
    },
    {
        "title": "The Modern Grid",
        "subtitle": "Things got complicated",
        "bullets": ["Electricity flows TWO ways — homes send power back", "Millions of devices connected simultaneously", "Needs real-time balancing every single second", "Too complex for humans to manage manually"],
        "funny": "Now things got seriously complicated",
        "visual": "Ultra clean 3D isometric infographic, blue and yellow theme, soft lighting, corporate professional, no text — complex smart grid with homes having solar panels, EVs, batteries, bidirectional arrows flowing in multiple directions between many nodes and substations",
        "morph": "MORPH: New nodes and connections grow outward, network becomes denser, cloud begins forming above",
        "notes": "Now imagine millions of solar panels, EVs, batteries — all sending and receiving power. The grid became incredibly complex. Ask yourself: how do you control all this?",
        "is_title": False
    },
    {
        "title": "What is Cloud Computing?",
        "subtitle": "The brain we always needed",
        "bullets": ["Storage — keeps massive amounts of data safely", "Compute — processes data at incredible speed", "Scalable — grows bigger when you need more power", "AI execution — runs smart algorithms 24/7"],
        "funny": "Cloud = the brain we always needed",
        "visual": "Ultra clean 3D isometric infographic, blue and yellow theme, soft lighting, corporate professional, no text — large glowing cloud shape above a city, inside the cloud visible server racks with glowing blue and yellow lights, data streams flowing down to buildings below",
        "morph": "MORPH: Zoom into cloud — camera moves up into the cloud interior, servers become analytical screens",
        "notes": "Cloud computing is not just storage. It's a complete brain. It stores data, processes it, runs AI, and scales automatically. Think of it as renting a superbrain.",
        "is_title": False
    },
    {
        "title": "Role of Cloud in Energy",
        "subtitle": "",
        "bullets": ["Forecast demand — predict energy needs for tomorrow", "Predict renewable output — estimate solar and wind", "Control grid — automatically balance supply and demand", "Process millions of data points every second"],
        "funny": "Cloud thinks ahead so we don't panic",
        "visual": "Ultra clean 3D isometric infographic, blue and yellow theme, soft lighting, corporate professional, no text — cloud shape with internal screens showing line graphs, bar charts, and wave patterns, data flowing in from bottom, decisions flowing out as glowing arrows",
        "morph": "MORPH: Graphs inside cloud extend and grow, breaking out to show different patterns",
        "notes": "In energy, the cloud does three critical things. It forecasts what we'll need. It predicts what nature will give us. And it automatically balances the two. Every second.",
        "is_title": False
    },
    {
        "title": "The Challenges We Face",
        "subtitle": "",
        "bullets": ["Intermittency — sun hides, wind stops suddenly", "Congestion — too much power on one transmission line", "Data silos — each company keeps data to itself", "Legacy systems can't talk to modern devices"],
        "funny": "Too many cooks, zero communication",
        "visual": "Ultra clean 3D isometric infographic, blue and yellow theme, soft lighting, corporate professional, no text — chaotic grid scene with overloaded transmission lines glowing red, disconnected data blocks floating separately, solar panels partially shaded, wind turbines stopped",
        "morph": "MORPH: Chaos to organized — scattered elements smoothly slide into organized rows, red turns blue and yellow",
        "notes": "These are real daily problems for grid operators. Intermittency means renewables are unreliable. Congestion causes physical damage. Data silos mean no one sees the full picture.",
        "is_title": False
    },
    {
        "title": "The Cloud Solutions",
        "subtitle": "",
        "bullets": ["Monitoring — see everything across the grid 24/7", "Prediction — AI tells you what happens before it happens", "Optimization — automatically find best power distribution", "Response time drops from hours to milliseconds"],
        "funny": "See -> Think -> Act. Simple as that",
        "visual": "Ultra clean 3D isometric infographic, blue and yellow theme, soft lighting, corporate professional, no text — three connected steps as 3D platforms: first with eye icon and sensors, second with brain and AI patterns, third with gears and output arrows, connected by flowing data streams",
        "morph": "MORPH: Three platforms align horizontally then stack vertically into architecture layers",
        "notes": "The cloud solves each challenge systematically. Monitor everything. Predict what's coming. Optimize automatically. What used to take hours now takes milliseconds.",
        "is_title": False
    },
    {
        "title": "System Architecture",
        "subtitle": "Three layers working together",
        "bullets": ["Edge Layer — sensors and devices at grid site (fast response)", "Communication Layer — 5G, fiber, IoT networks (data highway)", "Cloud Layer — AI, storage, control center (the big brain)", "All three layers work together continuously"],
        "funny": "Full system pipeline — no weak links",
        "visual": "Ultra clean 3D isometric infographic, blue and yellow theme, soft lighting, corporate professional, no text — three horizontal layers stacked: bottom with small sensor devices on grid equipment, middle with glowing fiber optic lines and 5G towers, top with large cloud containing AI brain",
        "morph": "MORPH: Zoom into bottom edge layer, sensors enlarge to fill frame, then split into comparison",
        "notes": "Think of it like your body. Edge = your nerves and reflexes. Communication = your spinal cord. Cloud = your brain. All three must work together.",
        "is_title": False
    },
    {
        "title": "Edge vs Cloud",
        "subtitle": "Reflex vs Brain",
        "bullets": ["Edge — reacts in under 100 milliseconds (like a reflex)", "Cloud — takes 1-2 seconds but makes smart decisions", "Edge handles emergencies, Cloud handles strategy", "Example: Edge stops short circuit, Cloud plans tomorrow's energy"],
        "funny": "Reflex vs Brain — both needed to survive",
        "visual": "Ultra clean 3D isometric infographic, blue and yellow theme, soft lighting, corporate professional, no text — split comparison, left side shows small edge device on power pole with lightning-fast reflex lines, right side shows large cloud with deep thinking patterns",
        "morph": "MORPH: Left edge device glows brighter and moves to center, lines extend showing data flowing to cloud",
        "notes": "Why not put everything in the cloud? Because some things can't wait. A short circuit needs response in 10 milliseconds. That's the edge. But planning tomorrow's energy mix? That's the cloud.",
        "is_title": False
    },
    {
        "title": "Data Pipeline",
        "subtitle": "From sensors to action",
        "bullets": ["Collect — 50,000+ sensors send data every second", "Transmit — secure channels move data to cloud", "Process — AI cleans, analyzes, and finds patterns", "Act — commands sent back to grid equipment automatically"],
        "funny": "Data is the new electricity",
        "visual": "Ultra clean 3D isometric infographic, blue and yellow theme, soft lighting, corporate professional, no text — flowing pipeline left to right: small sensor nodes into communication tunnel, into cloud processing center, output arrows to grid equipment, continuous data particles flowing",
        "morph": "MORPH: Output arrows extend and branch out to four application icons that grow from endpoints",
        "notes": "This pipeline runs nonstop. 50,000 sensors, every second, all day. The data flows up, AI processes it, decisions flow back down. It's the nervous system of the grid.",
        "is_title": False
    },
    {
        "title": "Applications in Energy",
        "subtitle": "One cloud, many uses",
        "bullets": ["Solar farms — cloud predicts output, optimizes panel angles", "Wind farms — forecasts wind, schedules maintenance", "EV charging — balances load across thousands of stations", "Battery storage — decides when to store and when to release"],
        "funny": "One cloud to rule them all",
        "visual": "Ultra clean 3D isometric infographic, blue and yellow theme, soft lighting, corporate professional, no text — central cloud connected to four surrounding scenes: solar farm with tracking panels, wind farm with turbines, EV charging station with cars, battery storage facility",
        "morph": "MORPH: Four application nodes grow larger, central cloud pulses, AI brain patterns become visible",
        "notes": "The beauty of cloud is one platform serves many applications. Solar prediction, wind forecasting, EV management, battery control — all from the same cloud brain.",
        "is_title": False
    },
    {
        "title": "AI in Energy Systems",
        "subtitle": "Numbers that matter",
        "bullets": ["Forecasting accuracy — less than 2% error in prediction", "Predictive maintenance — finds failures 2 weeks before they happen", "Optimization — reduces energy waste by 15-20%", "Google DeepMind cut cooling energy by 40% in data centers"],
        "funny": "AI predicts better than our gut feeling",
        "visual": "Ultra clean 3D isometric infographic, blue and yellow theme, soft lighting, corporate professional, no text — large AI brain shape made of neural network nodes, connected to energy graphs showing accurate prediction lines overlapping actual usage, small equipment icons with green checkmarks",
        "morph": "MORPH: Prediction line extends forward on graph showing future predictions, transforms into holographic twin city",
        "notes": "These are real numbers. Less than 2% error is incredible when you're managing gigawatts. Finding failures 2 weeks early saves millions. And Google's 40% reduction is a proven case.",
        "is_title": False
    },
    {
        "title": "Digital Twin Technology",
        "subtitle": "Test before you touch",
        "bullets": ["Virtual copy — exact 3D digital model of the entire grid", "Simulation — test scenarios without touching real equipment", "What-if analysis — What if a storm hits this area?", "Utility companies save $1-3 billion annually with digital twins"],
        "funny": "Test before you mess it up",
        "visual": "Ultra clean 3D isometric infographic, blue and yellow theme, soft lighting, corporate professional, no text — real city grid on left in solid colors, holographic transparent twin city on right with glowing wireframe overlay, simulation particles flowing between them, storm cloud testing above twin",
        "morph": "MORPH: Holographic twin city rotates, protective shield layer appears around it transitioning to security",
        "notes": "A digital twin is like a video game version of your real grid. You can simulate storms, equipment failures, demand spikes — all without risking the real grid. It saves billions.",
        "is_title": False
    },
    {
        "title": "Security in Cloud Energy",
        "subtitle": "",
        "bullets": ["Encryption — all data locked with military-grade codes", "Access control — only authorized people enter the system", "Threat detection — AI spots hackers in real-time", "Energy sector faces 10,000+ cyber attacks every day globally"],
        "funny": "Hackers? Not today, not ever",
        "visual": "Ultra clean 3D isometric infographic, blue and yellow theme, soft lighting, corporate professional, no text — smart grid wrapped in glowing hexagonal shield pattern, lock icons at access points, red threat particles bouncing off shield, surveillance sensors scanning perimeter",
        "morph": "MORPH: Hexagonal shield expands and transforms into network connecting multiple homes",
        "notes": "10,000 attacks per day. That's not a typo. Energy infrastructure is the number 1 target for hackers. Cloud security with AI threat detection is no longer optional — it's survival.",
        "is_title": False
    },
    {
        "title": "Virtual Power Plant (VPP)",
        "subtitle": "1000 homes = 1 power plant",
        "bullets": ["Connect 1,000 home batteries as one virtual power plant", "Capacity: 1,000 homes x 10 kWh = 10 MWh virtual plant", "Homeowners earn money by sharing stored energy", "Tesla VPP in South Australia: 4,000 homes, 250 MWh"],
        "funny": "Together we are powerful",
        "visual": "Ultra clean 3D isometric infographic, blue and yellow theme, soft lighting, corporate professional, no text — many individual homes with small battery icons, glowing connection lines linking them together, merging into one large unified power plant shape on right side, energy particles flowing",
        "morph": "MORPH: Homes slide together and merge into power plant shape, then one home separates and EV drives out",
        "notes": "This is mind-blowing. 4,000 homes in Australia became a 250 MWh power plant — without building anything new. Just software connecting existing batteries. Cloud makes this possible.",
        "is_title": False
    },
    {
        "title": "Vehicle-to-Grid (V2G)",
        "subtitle": "Your car is a power bank",
        "bullets": ["Your EV battery stores energy and sends it back to grid", "One EV battery = 60-100 kWh of mobile storage", "10 million EVs = 600 GWh of grid storage", "Nissan Leaf V2G trial in UK: owners earned 1,500 pounds/year"],
        "funny": "Your car is now a power bank on wheels",
        "visual": "Ultra clean 3D isometric infographic, blue and yellow theme, soft lighting, corporate professional, no text — electric vehicle at home charging station with bidirectional glowing arrows, one arrow from grid to car showing charging, another from car to home showing discharging, battery indicators",
        "morph": "MORPH: Discharge arrow grows larger and extends to connect with data processing scene for case study",
        "notes": "Your car sits parked 95% of the time. During that time, its battery can store solar energy and feed it back during peak hours. Nissan proved owners can earn money doing this.",
        "is_title": False
    },
    {
        "title": "Case Study: IBM Global Services",
        "subtitle": "Making 50 billion rows useful",
        "bullets": ["Problem: 50 billion rows of energy data sitting unused", "Solution: moved to cloud with AI analytics platform", "Result: AI model training became 75% faster", "Impact: real-time grid optimization across 5 million smart meters"],
        "funny": "50 billion rows finally doing something useful",
        "visual": "Ultra clean 3D isometric infographic, blue and yellow theme, soft lighting, corporate professional, no text — left side shows massive messy data blocks piled up chaotically, right side shows same data organized into clean structured columns flowing into cloud with AI processing, speed indicators",
        "morph": "MORPH: Messy data blocks rearrange into clean columns, flow into cloud, outputs transform into trading charts",
        "notes": "50 billion. Let that number sink in. This is how much data a modern grid generates. IBM moved it to cloud and suddenly it became useful — 75% faster, 5 million meters optimized in real-time.",
        "is_title": False
    },
    {
        "title": "Case Study: Tesla Autobidder",
        "subtitle": "Energy trading = stock market",
        "bullets": ["Problem: battery storage sits idle when not needed", "Solution: AI cloud platform trades energy automatically", "Result: Tesla batteries in Australia earned $50 million in 2 years", "Scale: managing 3 GWh of battery storage across multiple sites"],
        "funny": "Energy trading = stock market for electricity",
        "visual": "Ultra clean 3D isometric infographic, blue and yellow theme, soft lighting, corporate professional, no text — large Tesla-style battery pack with glowing cells, connected to cloud showing trading interface with rising graph lines, money flow arrows, buy/sell indicators, multiple battery sites in network",
        "morph": "MORPH: Trading graph line rises continuously, grows off slide, transforms into branded cloud logo",
        "notes": "Tesla took batteries that were sitting idle and turned them into money-making machines. The AI buys energy when it's cheap and sells when it's expensive. $50 million in 2 years from batteries that were doing nothing.",
        "is_title": False
    },
    {
        "title": "Introducing EnerCloud Pvt Limited",
        "subtitle": "Cloud-native energy management from India",
        "bullets": ["Cloud-native energy management platform built in India", "Mission: make every watt of energy count", "Built for utilities, solar operators, EV networks, grid operators", "Core belief: software can solve energy's hardest problems"],
        "funny": "Our solution starts here",
        "visual": "Ultra clean 3D isometric infographic, blue and yellow theme, soft lighting, corporate professional, no text — branded cloud platform floating above Indian cityscape, glowing energy nodes connecting to cloud, modern data center elements, professional corporate feel",
        "morph": "MORPH: Camera pushes into EnerCloud platform, outer city fades, internal problem indicators appear",
        "notes": "Now let me introduce EnerCloud. We are building exactly what we've been talking about — a cloud platform that makes energy intelligent. Built right here in India, for India and the world.",
        "is_title": False
    },
    {
        "title": "The Problem EnerCloud Solves",
        "subtitle": "",
        "bullets": ["Imbalance: solar overproduction at noon crashes grid prices", "Waste: 30-40% of renewable energy gets curtailed (thrown away)", "Overload: evening peak causes transformer failures and blackouts", "No visibility: operators can't see what's happening across the grid"],
        "funny": "The problem is very real and very expensive",
        "visual": "Ultra clean 3D isometric infographic, blue and yellow theme, soft lighting, corporate professional, no text — before and after comparison, left shows red warning indicators on grid with wasted energy particles, overloaded transformers, right shows balanced green grid with optimized flow",
        "morph": "MORPH: Red before side shrinks and slides left, green after side expands to fill frame, transforms into flow diagram",
        "notes": "30 to 40 percent of renewable energy is literally thrown away because the grid can't handle it. That's like filling your water glass and pouring half on the floor. EnerCloud stops this waste.",
        "is_title": False
    },
    {
        "title": "How EnerCloud Works",
        "subtitle": "Three steps, one loop",
        "bullets": ["Step 1: Collect — gather data from 10,000+ grid endpoints every second", "Step 2: Analyze — AI processes data and predicts 24 hours ahead", "Step 3: Control — automated commands balance grid in real-time", "This cycle runs continuously, 24 hours, 365 days"],
        "funny": "Simple logic: Data -> AI -> Action",
        "visual": "Ultra clean 3D isometric infographic, blue and yellow theme, soft lighting, corporate professional, no text — three-step horizontal flow with glowing circular nodes, first shows sensor data streams, second shows AI brain with prediction waveforms, third shows grid control with balanced energy arrows",
        "morph": "MORPH: Data particles accelerate through pipeline, three nodes transform into dashboard interface panels",
        "notes": "It sounds simple because it is — at the concept level. Collect data. Let AI analyze it. Automatically control the grid. The complexity is hidden inside the AI. For the operator, it's three steps.",
        "is_title": False
    },
    {
        "title": "EnerCloud Key Features",
        "subtitle": "Everything in one platform",
        "bullets": ["Live Monitoring — see every watt flowing in real-time", "Smart Prediction — 24-hour forecast with 97% accuracy", "Auto Optimization — reduces peak load by 20-30% automatically", "One Dashboard — control solar, wind, EV, battery from one screen"],
        "funny": "Everything in one place — finally",
        "visual": "Ultra clean 3D isometric infographic, blue and yellow theme, soft lighting, corporate professional, no text — modern 3D dashboard interface with multiple floating panels showing line graphs, pie charts, map views, gauge meters, all in blue and yellow, data points glowing",
        "morph": "MORPH: Camera zooms into map panel on dashboard, map expands to show full city with balanced energy flow",
        "notes": "97% accuracy on 24-hour predictions. 20-30% peak reduction without human intervention. One screen to control everything. These aren't promises — these are our current metrics.",
        "is_title": False
    },
    {
        "title": "EnerCloud Real Use Case",
        "subtitle": "Proven results",
        "bullets": ["Deployed across 10,000 homes in a smart city project", "Evening peak demand reduced by 25%", "Utility company saved 12 crore rupees annually", "Zero blackouts during summer peak for the first time"],
        "funny": "No blackout = happy everyone",
        "visual": "Ultra clean 3D isometric infographic, blue and yellow theme, soft lighting, corporate professional, no text — balanced smart city at evening with homes glowing warmly, EnerCloud platform visible above with balanced distribution, green checkmarks on substations, smooth energy flow lines",
        "morph": "MORPH: Balanced city energy flow lines extend outward and upward, city pulls back to reveal wider futuristic landscape",
        "notes": "This is not a simulation. 10,000 real homes. 25% peak reduction. 12 crore saved. And for the first time in that city's history — zero blackouts during summer. This is what cloud plus energy looks like.",
        "is_title": False
    },
    {
        "title": "The Future is Now",
        "subtitle": "Energy + Intelligence = The Future",
        "bullets": ["Every grid will be cloud-connected and AI-driven", "Software will control physical energy infrastructure", "EnerCloud is building that future starting today", "The grid of 2030 thinks, learns, and heals itself"],
        "funny": "Energy + Intelligence = The Future",
        "visual": "Ultra clean 3D isometric infographic, blue and yellow theme, soft lighting, corporate professional, no text — futuristic city skyline at golden hour, massive glowing cloud brain above with neural connections to every building, solar panels, wind turbines, EVs, all interconnected, optimistic inspiring atmosphere",
        "morph": "MORPH: Camera slowly pulls back revealing wider futuristic landscape, energy particles fade into stars, fade to black",
        "notes": "The grid of the future won't need humans to operate it. It will think. It will learn. It will heal itself when something breaks. That future isn't 20 years away. EnerCloud is building it today. Thank you.",
        "is_title": False
    }
]

def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, italic=False, color=DARK_TEXT, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.alignment = align
    return txBox

def add_shape(slide, shape_type, left, top, width, height, fill_color=None, line_visible=False):
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    shape.line.fill.background()
    return shape

for idx, sd in enumerate(slides_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    slide_num = idx + 1

    if sd["is_title"]:
        # Dark background
        add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 7.5, DARK_BLUE)

        # Decorative circles
        c1 = add_shape(slide, MSO_SHAPE.OVAL, 10.5, -1, 4, 4, MED_BLUE)
        c1.fill.fore_color.rgb = MED_BLUE
        # Set transparency via XML workaround
        from pptx.oxml.ns import qn
        solidFill = c1.fill._fill
        srgb = solidFill.find(qn('a:solidFill')).find(qn('a:srgbClr'))
        alpha = srgb.makeelement(qn('a:alpha'), {'val': '15000'})
        srgb.append(alpha)

        c2 = add_shape(slide, MSO_SHAPE.OVAL, -1.5, 5.5, 3.5, 3.5, YELLOW)
        solidFill2 = c2.fill._fill
        srgb2 = solidFill2.find(qn('a:solidFill')).find(qn('a:srgbClr'))
        alpha2 = srgb2.makeelement(qn('a:alpha'), {'val': '10000'})
        srgb2.append(alpha2)

        # Yellow accent line
        add_shape(slide, MSO_SHAPE.RECTANGLE, 1, 2.8, 0.8, 0.06, YELLOW)

        # Title
        add_textbox(slide, 1, 3.0, 11, 1.8, sd["title"], 42, True, False, WHITE)

        # Subtitle
        add_textbox(slide, 1, 4.6, 11, 0.6, sd["subtitle"], 18, False, False, LIGHT_BLUE)

        # Funny line
        add_textbox(slide, 1, 5.5, 11, 0.5, sd["funny"], 15, False, True, YELLOW)

        # Cloud icon
        cloud = add_shape(slide, MSO_SHAPE.CLOUD, 10.2, 6.0, 2, 1.2, MED_BLUE)
        solidFill3 = cloud.fill._fill
        srgb3 = solidFill3.find(qn('a:solidFill')).find(qn('a:srgbClr'))
        alpha3 = srgb3.makeelement(qn('a:alpha'), {'val': '70000'})
        srgb3.append(alpha3)

    else:
        # White background
        add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 7.5, WHITE)

        # Top bar
        add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 0.08, DARK_BLUE)
        add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0.08, 13.333, 0.04, YELLOW)

        # Side bar
        add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 0.08, 7.5, MED_BLUE)

        # Slide number
        add_textbox(slide, 12.2, 0.25, 0.8, 0.3, f"{slide_num}/26", 10, False, False, GRAY, PP_ALIGN.RIGHT)

        # Title
        add_textbox(slide, 0.6, 0.4, 11.5, 0.7, sd["title"], 30, True, False, DARK_BLUE)

        # Subtitle
        if sd["subtitle"]:
            add_textbox(slide, 0.6, 1.05, 11.5, 0.4, sd["subtitle"], 14, False, True, MED_BLUE)

        # Yellow line
        add_shape(slide, MSO_SHAPE.RECTANGLE, 0.6, 1.45, 1, 0.04, YELLOW)

        # Bullets
        by = 1.7
        for b_idx, bullet in enumerate(sd["bullets"]):
            # Yellow dot
            dot = add_shape(slide, MSO_SHAPE.OVAL, 0.65, by + 0.08, 0.12, 0.12, YELLOW)
            # Text
            add_textbox(slide, 0.95, by - 0.02, 11.2, 0.55, bullet, 17, False, False, DARK_TEXT)
            by += 0.65

        # Funny line box
        fy = by + 0.25
        funny_box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.6, fy, 12, 0.5, DARK_BLUE)
        add_textbox(slide, 0.85, fy + 0.08, 11.5, 0.4, sd["funny"], 13, False, True, YELLOW)

        # Visual prompt area
        vis_y = 6.3
        add_shape(slide, MSO_SHAPE.RECTANGLE, 0.6, vis_y, 12, 0.8, LIGHT_BG)
        vis_text = sd["visual"]
        if len(vis_text) > 200:
            vis_text = vis_text[:200] + "..."
        add_textbox(slide, 0.8, vis_y + 0.08, 11.6, 0.7, f"[IMAGE: {vis_text}]", 8, False, False, GRAY)

        # Morph note
        morph_text = sd["morph"]
        if len(morph_text) > 120:
            morph_text = morph_text[:120] + "..."
        add_textbox(slide, 9, 7.15, 4.2, 0.3, morph_text, 7, False, False, RGBColor(0xBB, 0xBB, 0xBB), PP_ALIGN.RIGHT)

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = (
        f"SPEAKER NOTES:\n\n{sd['notes']}\n\n"
        f"--- VISUAL PROMPT ---\n{sd['visual']}\n\n"
        f"--- MORPH TRANSITION ---\n{sd['morph']}"
    )

# Save
output_path = "EnerCloud_Presentation.pptx"
prs.save(output_path)
print(f"DONE! File saved: {output_path}")
print(f"Total slides: {len(slides_data)}")